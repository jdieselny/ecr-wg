import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Brand Colors (Adapted from jdieselny genetics logo theme)
COLOR_MIDNIGHT = RGBColor(0x0C, 0x11, 0x17)  # #0C1117 (Deep slate background)
COLOR_ICE_BLUE = RGBColor(0x8E, 0xCA, 0xE6)  # #8ECAE6 (Ice Blue accent)
COLOR_SLATE = RGBColor(0x47, 0x55, 0x69)     # #475569 (Muted slate text/accents)
COLOR_CHARCOAL = RGBColor(0x0F, 0x17, 0x2A)  # #0F172A (Primary dark body text)
COLOR_LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)  # #F8FAFC (Light mode background)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)     # #FFFFFF (White text/fills)
COLOR_BORDER = RGBColor(0xE2, 0xE8, 0xF0)    # #E2E8F0 (Table/border line)

def create_full_background(slide, color):
    """Fills the slide background with a solid color."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()  # No border
    # Send background to back by re-ordering shapes
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg

def add_corner_branding(slide, text_color):
    """Adds subtle corner details (brackets and top-right text) from the brand logo."""
    # Top-Right Label: JD//NY - 001
    tx_box = slide.shapes.add_textbox(Inches(10.5), Inches(0.4), Inches(2.3), Inches(0.4))
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = "ECR-WG // REPLICATION - 04"
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = "Segoe UI"
    p.font.size = Pt(8.5)
    p.font.color.rgb = text_color
    p.font.bold = True
    
    # Corner brackets (Top-Left)
    line1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(0.3), Inches(0.015))
    line1.fill.solid()
    line1.fill.fore_color.rgb = text_color
    line1.line.fill.background()
    
    line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(0.015), Inches(0.3))
    line2.fill.solid()
    line2.fill.fore_color.rgb = text_color
    line2.line.fill.background()

    # Corner brackets (Bottom-Right)
    line3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12.5), Inches(6.985), Inches(0.3), Inches(0.015))
    line3.fill.solid()
    line3.fill.fore_color.rgb = text_color
    line3.line.fill.background()
    
    line4 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12.785), Inches(6.7), Inches(0.015), Inches(0.3))
    line4.fill.solid()
    line4.fill.fore_color.rgb = text_color
    line4.line.fill.background()

def add_horizontal_rule(slide, top_y, height_y=0.02, color=COLOR_ICE_BLUE):
    """Draws a horizontal accent line across the slide."""
    hr = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(top_y), Inches(11.73), Inches(height_y)
    )
    hr.fill.solid()
    hr.fill.fore_color.rgb = color
    hr.line.fill.background()

def add_slide_header(slide, title_text, dark_mode=False):
    """Adds a standardized slide title."""
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.73), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Segoe UI"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE if dark_mode else COLOR_MIDNIGHT
    
    # Add horizontal rule below title
    add_horizontal_rule(slide, top_y=1.5, color=COLOR_ICE_BLUE)

# Initialize Presentation and configure widescreen 16:9 dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

# ==============================================================================
# SLIDE 1: Title Slide (Dark Mode)
# ==============================================================================
slide1 = prs.slides.add_slide(blank_layout)
create_full_background(slide1, COLOR_MIDNIGHT)
add_corner_branding(slide1, COLOR_SLATE)

# Title & Subtitle in one text frame
title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.73), Inches(3.5))
tf = title_box.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

p_title = tf.paragraphs[0]
p_title.text = "NETL ENERGY BENCHMARK"
p_title.font.name = "Segoe UI"
p_title.font.size = Pt(40)
p_title.font.bold = True
p_title.font.color.rgb = COLOR_WHITE

p_sub = tf.add_paragraph()
p_sub.text = "EXHIBIT D: Independent Replication // Gemini 2.5 Flash"
p_sub.space_before = Pt(8)
p_sub.font.name = "Segoe UI"
p_sub.font.size = Pt(18)
p_sub.font.color.rgb = COLOR_ICE_BLUE

# Decorative Divider
add_horizontal_rule(slide1, top_y=3.8, height_y=0.03, color=COLOR_ICE_BLUE)

# Metadata paragraph
p_meta = tf.add_paragraph()
p_meta.text = "EFFICIENCY-CENTERED REASONING WORKING GROUP (ECR-WG)\nRegistrant: Justin Kintzele  |  Date: 2026-05-22"
p_meta.space_before = Pt(40)
p_meta.font.name = "Segoe UI"
p_meta.font.size = Pt(11)
p_meta.font.color.rgb = COLOR_SLATE

# ==============================================================================
# SLIDE 2: Replication Verdict (Light Mode)
# ==============================================================================
slide2 = prs.slides.add_slide(blank_layout)
create_full_background(slide2, COLOR_LIGHT_BG)
add_slide_header(slide2, "EXHIBIT D // Replication Verdict")

# Content block
content_box = slide2.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.73), Inches(4.5))
tf = content_box.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

def add_bullet(tf, title, body):
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.space_after = Pt(14)
    run_title = p.add_run()
    run_title.text = "▪  " + title + ": "
    run_title.font.bold = True
    run_title.font.size = Pt(14)
    run_title.font.color.rgb = COLOR_MIDNIGHT
    
    run_body = p.add_run()
    run_body.text = body
    run_body.font.size = Pt(14)
    run_body.font.color.rgb = COLOR_CHARCOAL
    p.font.name = "Segoe UI"

add_bullet(tf, "Objective", "Re-run Exhibit A's Flash cell on later harness commit (0736720) with a separate orchestrator node (Node B) to verify savings durability.")
add_bullet(tf, "Key Finding", "All four headline economics metrics replicate within < 2 percentage points of the original Node A benchmarks. Savings are confirmed to be structurally inherent, not an artifact of specific runs.")
add_bullet(tf, "Task Quality", "Demonstrates a +6.7% improvement in task quality (0.914 vs 0.857) alongside radical compute compression.")
add_bullet(tf, "Significance", "Replication provides empirical grounding (N=5 per arm) verifying that ECR-WG overlays deliver consistent, portable, and reliable efficiency gains.")

# ==============================================================================
# SLIDE 3: Token Economics Table (Light Mode)
# ==============================================================================
slide3 = prs.slides.add_slide(blank_layout)
create_full_background(slide3, COLOR_LIGHT_BG)
add_slide_header(slide3, "Gemini 2.5 Flash Token Economics (N=5)")

# Create Table
rows, cols = 6, 5
table_shape = slide3.shapes.add_table(rows, cols, Inches(0.8), Inches(2.0), Inches(11.73), Inches(4.5))
table = table_shape.table

# Column widths
table.columns[0].width = Inches(3.23)
table.columns[1].width = Inches(2.2)
table.columns[2].width = Inches(2.2)
table.columns[3].width = Inches(2.1)
table.columns[4].width = Inches(2.0)

headers = ["Metric", "Arm A (Stateless)", "Arm B (Orchestrated)", "Delta", "Savings %"]
data = [
    ["Input Tokens (Prefill)", "86,321 +/- 8,475", "28,224 +/- 2,316", "-58,097", "67.3% Saved"],
    ["Input Charged Prefill", "21,648 +/- 1,797", "13,615 +/- 8,560", "-8,033", "37.1% Saved"],
    ["Output Tokens (Gen)", "16,999 +/- 1,648", "3,978 +/- 175", "-13,021", "76.6% Saved"],
    ["Wall-Clock Time", "112.2s +/- 8.0s", "37.0s +/- 3.4s", "-75.2s", "67.0% Saved"],
    ["Task Quality (mean)", "0.857", "0.914", "+0.057", "+6.7% Gain"]
]

# Write header
for c, header_text in enumerate(headers):
    cell = table.cell(0, c)
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_MIDNIGHT
    p = cell.text_frame.paragraphs[0]
    p.text = header_text
    p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
    p.font.name = "Segoe UI"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

# Write data
for r, row_data in enumerate(data):
    for c, val in enumerate(row_data):
        cell = table.cell(r + 1, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_WHITE if r % 2 == 0 else COLOR_LIGHT_BG
        p = cell.text_frame.paragraphs[0]
        p.text = val
        p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.bold = (c == 4)  # Bold savings column
        
        # Color specific elements
        if c == 4:
            p.font.color.rgb = RGBColor(0x16, 0xA3, 0x4A) if "Saved" in val or "Gain" in val else COLOR_CHARCOAL
        else:
            p.font.color.rgb = COLOR_CHARCOAL

# ==============================================================================
# SLIDE 4: Cache Behavior & CAPEX (Light Mode)
# ==============================================================================
slide4 = prs.slides.add_slide(blank_layout)
create_full_background(slide4, COLOR_LIGHT_BG)
add_slide_header(slide4, "Cache Dynamics & CAPEX Implication")

# Callout Shaded Box
shading = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.0), Inches(11.73), Inches(2.2))
shading.fill.solid()
shading.fill.fore_color.rgb = COLOR_WHITE
shading.line.color.rgb = COLOR_BORDER

# Shading left accent bar
accent_bar = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.0), Inches(0.1), Inches(2.2))
accent_bar.fill.solid()
accent_bar.fill.fore_color.rgb = COLOR_ICE_BLUE
accent_bar.line.fill.background()

# Callout text box
callout_box = slide4.shapes.add_textbox(Inches(1.1), Inches(2.15), Inches(11.1), Inches(1.9))
tf = callout_box.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
p = tf.paragraphs[0]
p.text = "LOAD-BEARING ANALYSIS: Automatic Prompt Cache vs. Orchestration"
p.font.bold = True
p.font.size = Pt(13)
p.font.color.rgb = COLOR_MIDNIGHT
p.font.name = "Segoe UI"

p2 = tf.add_paragraph()
p2.text = (
    "In this cell, Gemini's automatic prompt cache absorbed a larger portion of the stateless arm's input. "
    "This compressed the \"input charged\" prefill delta to 37.1% (compared to the 67.3% total input savings).\n\n"
    "Because peak capacity demands dictate infrastructure CAPEX requirements, the \"input charged\" metric—reflecting "
    "what actually hits physical GPU cores—is the true load-bearing argument for federal proposals, not raw input savings."
)
p2.space_before = Pt(8)
p2.font.size = Pt(12)
p2.font.color.rgb = COLOR_SLATE
p2.font.italic = True
p2.font.name = "Segoe UI"

# Additional bullets below callout
bullets_box = slide4.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.73), Inches(2.2))
tf2 = bullets_box.text_frame
tf2.word_wrap = True
tf2.margin_left = tf2.margin_right = tf2.margin_top = tf2.margin_bottom = 0

add_bullet(tf2, "CAPEX Deferral", "Real-world infrastructure scaling limits are bound to uncached prefill costs. Overlays defer these capital expenditures directly by intercepting redundant requests before they reach the model provider.")
add_bullet(tf2, "Provider Independence", "Relying on automatic, opaque provider-side caching leaves cost structures volatile. ECR-WG overlays guarantee deterministic memory control at the boundary plane.")

# ==============================================================================
# SLIDE 5: Recipe & Reproducibility (Dark Mode)
# ==============================================================================
slide5 = prs.slides.add_slide(blank_layout)
create_full_background(slide5, COLOR_MIDNIGHT)
add_corner_branding(slide5, COLOR_SLATE)
add_slide_header(slide5, "Exhibit Limitations & Reproduction", dark_mode=True)

# Main columns (using textboxes side-by-side)
# Left Column: Limitations
left_box = slide5.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5))
tf_left = left_box.text_frame
tf_left.word_wrap = True
tf_left.margin_left = tf_left.margin_right = tf_left.margin_top = tf_left.margin_bottom = 0

p_lh = tf_left.paragraphs[0]
p_lh.text = "EXHIBIT LIMITATIONS"
p_lh.font.bold = True
p_lh.font.size = Pt(13)
p_lh.font.color.rgb = COLOR_ICE_BLUE
p_lh.font.name = "Segoe UI"

def add_dark_bullet(tf, title, body):
    p = tf.add_paragraph()
    p.space_before = Pt(10)
    p.font.name = "Segoe UI"
    
    run_t = p.add_run()
    run_t.text = "▪  " + title + ": "
    run_t.font.bold = True
    run_t.font.size = Pt(12)
    run_t.font.color.rgb = COLOR_WHITE
    
    run_b = p.add_run()
    run_b.text = body
    run_b.font.size = Pt(12)
    run_b.font.color.rgb = COLOR_SLATE

add_dark_bullet(tf_left, "Sample Size (N=5)", "Sufficient for validating directionality, but not powered for high-confidence variance or tight effect-size estimations.")
add_dark_bullet(tf_left, "Single Scenario", "Currently limited to 'code-review-iteration-01'. Additional benchmark scenarios are required to broaden the evidence base.")
add_dark_bullet(tf_left, "Automatic Cache Volatility", "Gemini's caching mechanics are proprietary and opaque, introducing external noise into the charged token calculations.")

# Right Column: Recipe
right_box = slide5.shapes.add_textbox(Inches(6.8), Inches(2.0), Inches(5.7), Inches(4.5))
tf_right = right_box.text_frame
tf_right.word_wrap = True
tf_right.margin_left = tf_right.margin_right = tf_right.margin_top = tf_right.margin_bottom = 0

p_rh = tf_right.paragraphs[0]
p_rh.text = "REPRODUCTION RECIPE"
p_rh.font.bold = True
p_rh.font.size = Pt(13)
p_rh.font.color.rgb = COLOR_ICE_BLUE
p_rh.font.name = "Segoe UI"

p_desc = tf_right.add_paragraph()
p_desc.text = "Run the replicate commands from the repository root. Requires a validated PROVIDER_API_KEY inside the local environment:"
p_desc.space_before = Pt(10)
p_desc.font.size = Pt(12)
p_desc.font.color.rgb = COLOR_SLATE
p_desc.font.name = "Segoe UI"

# Command box (Shaded background for code block in PPTX)
cmd_bg = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(3.2), Inches(5.7), Inches(2.2))
cmd_bg.fill.solid()
cmd_bg.fill.fore_color.rgb = RGBColor(0x1B, 0x22, 0x2C) # Lighter slate block background
cmd_bg.line.fill.background()

cmd_box = slide5.shapes.add_textbox(Inches(6.9), Inches(3.3), Inches(5.5), Inches(2.0))
tf_cmd = cmd_box.text_frame
tf_cmd.word_wrap = True
p_cmd = tf_cmd.paragraphs[0]
p_cmd.text = (
    "python -m benchmarks.scripts.runner energy \\\n"
    "  --scenario code-review-iteration-01 --arm A_naive --n 5\n\n"
    "python -m benchmarks.scripts.runner energy \\\n"
    "  --scenario code-review-iteration-01 --arm B_engineered --n 5\n\n"
    "python -m benchmarks.scripts.analyzer \\\n"
    "  --results-dir benchmarks/results/energy/code-review-iteration-01"
)
p_cmd.font.name = "Consolas"
p_cmd.font.size = Pt(9)
p_cmd.font.color.rgb = COLOR_ICE_BLUE

# Save presentation
output_pptx = "evidence/netl-energy-v0.4.pptx"
prs.save(output_pptx)
print(f"Successfully generated presentation slide deck: {output_pptx}")
