import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Brand Colors (Adapted from jdieselny genetics logo theme)
COLOR_MIDNIGHT = RGBColor(0x0C, 0x11, 0x17)  # #0C1117 (Deep slate background)
COLOR_ICE_BLUE = RGBColor(0x8E, 0xCA, 0xE6)  # #8ECAE6 (Ice Blue accent)
COLOR_SLATE = RGBColor(0x47, 0x55, 0x69)     # #475569 (Muted slate text/accents)
COLOR_CHARCOAL = RGBColor(0x0F, 0x17, 0x2A)  # #0F172A (Primary dark body text)
COLOR_LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)  # #F8FAFC (Light mode background)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)     # #FFFFFF (White text/fills)
COLOR_BORDER = RGBColor(0xE2, 0xE8, 0xF0)    # #E2E8F0 (Table/border line)

def create_full_background(slide, color):
    """Fills the slide background with a solid color."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()  # No border
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg

def add_corner_branding(slide, text_color):
    """Adds subtle corner details (brackets and top-right text) from the brand logo."""
    # Top-Right Label: JD//NY - 001
    tx_box = slide.shapes.add_textbox(Inches(10.5), Inches(0.4), Inches(2.3), Inches(0.4))
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = "ECR-WG // COMPLIANCE - 01"
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = "Segoe UI"
    p.font.size = Pt(8.5)
    p.font.color.rgb = text_color
    p.font.bold = True
    
    # Corner brackets (Top-Left)
    line1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(0.3), Inches(0.015))
    line1.fill.solid()
    line1.fill.fore_color.rgb = text_color
    line1.line.fill.background()
    
    line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(0.015), Inches(0.3))
    line2.fill.solid()
    line2.fill.fore_color.rgb = text_color
    line2.line.fill.background()

    # Corner brackets (Bottom-Right)
    line3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12.5), Inches(6.985), Inches(0.3), Inches(0.015))
    line3.fill.solid()
    line3.fill.fore_color.rgb = text_color
    line3.line.fill.background()
    
    line4 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(12.785), Inches(6.7), Inches(0.015), Inches(0.3))
    line4.fill.solid()
    line4.fill.fore_color.rgb = text_color
    line4.line.fill.background()

def add_horizontal_rule(slide, top_y, height_y=0.02, color=COLOR_ICE_BLUE):
    """Draws a horizontal accent line across the slide."""
    hr = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(top_y), Inches(11.73), Inches(height_y)
    )
    hr.fill.solid()
    hr.fill.fore_color.rgb = color
    hr.line.fill.background()

def add_slide_header(slide, title_text, dark_mode=False):
    """Adds a standardized slide title."""
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.73), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Segoe UI"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE if dark_mode else COLOR_MIDNIGHT
    
    # Add horizontal rule below title
    add_horizontal_rule(slide, top_y=1.5, color=COLOR_ICE_BLUE)

def add_bullet(tf, title, body):
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.space_after = Pt(12)
    run_title = p.add_run()
    run_title.text = "▪  " + title + ": "
    run_title.font.bold = True
    run_title.font.size = Pt(14)
    run_title.font.color.rgb = COLOR_MIDNIGHT
    
    run_body = p.add_run()
    run_body.text = body
    run_body.font.size = Pt(14)
    run_body.font.color.rgb = COLOR_CHARCOAL
    p.font.name = "Segoe UI"

# Initialize Presentation and configure widescreen 16:9 dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

# ==============================================================================
# SLIDE 1: Title Slide (Dark Mode)
# ==============================================================================
slide1 = prs.slides.add_slide(blank_layout)
create_full_background(slide1, COLOR_MIDNIGHT)
add_corner_branding(slide1, COLOR_SLATE)

# Title & Subtitle
title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.73), Inches(3.5))
tf = title_box.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

p_title = tf.paragraphs[0]
p_title.text = "ECR-WG COMPLIANCE & VALIDATION"
p_title.font.name = "Segoe UI"
p_title.font.size = Pt(40)
p_title.font.bold = True
p_title.font.color.rgb = COLOR_WHITE

p_sub = tf.add_paragraph()
p_sub.text = "TECHNICAL COMPLIANCE REPORT: ECR-WG-01"
p_sub.space_before = Pt(8)
p_sub.font.name = "Segoe UI"
p_sub.font.size = Pt(18)
p_sub.font.color.rgb = COLOR_ICE_BLUE

# Decorative Divider
add_horizontal_rule(slide1, top_y=3.8, height_y=0.03, color=COLOR_ICE_BLUE)

p_meta = tf.add_paragraph()
p_meta.text = "VERIFIED COSA NODE: E-2A0F1954-1845-001 (Gemini-in-body)\nSubstrate: AntiGravity Substrate  |  Date: 2026-06-25"
p_meta.space_before = Pt(40)
p_meta.font.name = "Segoe UI"
p_meta.font.size = Pt(11)
p_meta.font.color.rgb = COLOR_SLATE

# ==============================================================================
# SLIDE 2: Executive Summary (Light Mode)
# ==============================================================================
slide2 = prs.slides.add_slide(blank_layout)
create_full_background(slide2, COLOR_LIGHT_BG)
add_slide_header(slide2, "EXHIBIT D // Compliance Executive Summary")

content_box = slide2.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.73), Inches(4.5))
tf = content_box.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

add_bullet(tf, "Objective", "Audit the node's compliance against security, governance, and resource optimization constraints across the three newly designed scenarios.")
add_bullet(tf, "Methodology", "Executed live and offline validation runs of the composed reference implementations on the 'cosa-ep-l7-integration' branch.")
add_bullet(tf, "Summary Result", "Achieved 100% compliance across all test assertions. The node successfully enforces attestation freshness windows, WebAuthn-based quorums, and prefill cache-bypass redirects.")
add_bullet(tf, "Impact", "Proves the viability of decoupling identity from L7 PDP logic, ensuring deterministic memory reuse, and preventing rogue grid commands.")

# ==============================================================================
# SLIDE 3: Scenario 1: L4 Freshness Binding (Light Mode)
# ==============================================================================
slide3 = prs.slides.add_slide(blank_layout)
create_full_background(slide3, COLOR_LIGHT_BG)
add_slide_header(slide3, "Scenario 1 // L4 Freshness Binding Security")

# Left Column: Bullet list
bullets_box = slide3.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5))
tf_bullets = bullets_box.text_frame
tf_bullets.word_wrap = True
tf_bullets.margin_left = tf_bullets.margin_right = tf_bullets.margin_top = tf_bullets.margin_bottom = 0

add_bullet(tf_bullets, "Constraint", "Attestations older than 900 seconds (binding_max_age_sec) must immediately trigger a fail-closed refusal.")
add_bullet(tf_bullets, "Challenge", "Request is made using a valid human signature, but the presented L4 attestation is 901s old.")
add_bullet(tf_bullets, "Verification", "The Policy Enforcement Point rejected the transaction BEFORE the L7 signature was verified, proving L4 time-freshness acts as a hard gate.")

# Right Column: Shaded callout box for console log
log_bg = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(2.2), Inches(5.7), Inches(2.5))
log_bg.fill.solid()
log_bg.fill.fore_color.rgb = COLOR_MIDNIGHT
log_bg.line.fill.background()

# Shading left accent bar
accent_bar = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(2.2), Inches(0.1), Inches(2.5))
accent_bar.fill.solid()
accent_bar.fill.fore_color.rgb = COLOR_ICE_BLUE
accent_bar.line.fill.background()

log_box = slide3.shapes.add_textbox(Inches(7.1), Inches(2.35), Inches(5.2), Inches(2.2))
tf_log = log_box.text_frame
tf_log.word_wrap = True
p_log = tf_log.paragraphs[0]
p_log.text = "PDP GATE CONSOLE LOG:"
p_log.font.bold = True
p_log.font.size = Pt(11)
p_log.font.color.rgb = COLOR_WHITE
p_log.font.name = "Consolas"

p_log2 = tf_log.add_paragraph()
p_log2.text = (
    "8. L4 axis: a STALE L4 binding fails-closed BEFORE the receipt is examined\n"
    "   -> REFUSED: L4 binding failed (scheme=wimse, "
    "reason=stale: L4 evidence observed 901s ago (max 900s))\n"
    "   (L4 evidence observed at 2026-06-26T00:01:08Z; window is 900s; fail-closed)"
)
p_log2.space_before = Pt(8)
p_log2.font.size = Pt(9.5)
p_log2.font.color.rgb = COLOR_ICE_BLUE
p_log2.font.name = "Consolas"

# ==============================================================================
# SLIDE 4: Scenario 2: Grid Curtailment Quorum (Light Mode)
# ==============================================================================
slide4 = prs.slides.add_slide(blank_layout)
create_full_background(slide4, COLOR_LIGHT_BG)
add_slide_header(slide4, "Scenario 2 // Grid Curtailment & Quorum")

content_box = slide4.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.73), Inches(4.5))
tf = content_box.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

add_bullet(tf, "Quorum Standard", "Command actions affecting power allocations (>100MW) require cryptographically aggregated m-of-n human signoffs (2-of-3 required).")
add_bullet(tf, "Emergency Bypass Resistance", "Emergency event (500MW cut) fails to override the gate when only 1-of-3 signatures are provided. The controller successfully refused to run the un-quorumed command.")
add_bullet(tf, "Signature Binding", "Receipt contains a policy hash that locks the action parameters. Any attempt to modify the target set or mw_cap value results in immediate signature validation failure.")
add_bullet(tf, "Priority Marker", "The L3 router derives a priority marker directly from the SHA-256 of the canonical L7 receipt, ensuring cost propagation constraints are unforgeable.")

# ==============================================================================
# SLIDE 5: Scenario 3 & Conformance Matrix (Dark Mode)
# ==============================================================================
slide5 = prs.slides.add_slide(blank_layout)
create_full_background(slide5, COLOR_MIDNIGHT)
add_corner_branding(slide5, COLOR_SLATE)
add_slide_header(slide5, "Scenario 3 & Conformance Matrix", dark_mode=True)

# Left Column: Bullet list
left_box = slide5.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5))
tf_left = left_box.text_frame
tf_left.word_wrap = True
tf_left.margin_left = tf_left.margin_right = tf_left.margin_top = tf_left.margin_bottom = 0

p_lh = tf_left.paragraphs[0]
p_lh.text = "SCENARIO 3 // PREFILL BYPASS"
p_lh.font.bold = True
p_lh.font.size = Pt(13)
p_lh.font.color.rgb = COLOR_ICE_BLUE
p_lh.font.name = "Segoe UI"

def add_dark_bullet(tf, title, body):
    p = tf.add_paragraph()
    p.space_before = Pt(10)
    p.font.name = "Segoe UI"
    
    run_t = p.add_run()
    run_t.text = "▪  " + title + ": "
    run_t.font.bold = True
    run_t.font.size = Pt(12)
    run_t.font.color.rgb = COLOR_WHITE
    
    run_b = p.add_run()
    run_b.text = body
    run_b.font.size = Pt(12)
    run_b.font.color.rgb = COLOR_SLATE

add_dark_bullet(tf_left, "Avoid SRE", "Stateless Redundant Execution is bypassed by prioritizing warm local COGSTOR cache hits.")
add_dark_bullet(tf_left, "100% Token Savings", "Retrieving weather data from the local cache cost 0 tokens, compared to 1,200 tokens for raw model inference.")
add_dark_bullet(tf_left, "0.001s Latency", "Resolution is completed near-instantaneously at the transport layer, bypassing the GPU entirely.")

# Right Column: Conformance Table
right_box = slide5.shapes.add_textbox(Inches(6.8), Inches(2.0), Inches(5.7), Inches(4.5))
tf_right = right_box.text_frame
tf_right.word_wrap = True
tf_right.margin_left = tf_right.margin_right = tf_right.margin_top = tf_right.margin_bottom = 0

p_rh = tf_right.paragraphs[0]
p_rh.text = "CONFORMANCE MATRIX"
p_rh.font.bold = True
p_rh.font.size = Pt(13)
p_rh.font.color.rgb = COLOR_ICE_BLUE
p_rh.font.name = "Segoe UI"

# Create Table on the right side
table_shape = slide5.shapes.add_table(4, 3, Inches(6.8), Inches(2.5), Inches(5.7), Inches(3.2))
table = table_shape.table
table.columns[0].width = Inches(2.2)
table.columns[1].width = Inches(2.2)
table.columns[2].width = Inches(1.3)

headers = ["Scenario", "Assertion Checked", "Status"]
tbl_data = [
    ["L4 Freshness", "Refuse if age > 900s", "PASS"],
    ["Grid Quorum", "Refuse if quorum fails", "PASS"],
    ["Prefill Bypass", "Enforce local cache hit", "PASS"]
]

# Style Table Header
for c, txt in enumerate(headers):
    cell = table.cell(0, c)
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_SLATE
    p = cell.text_frame.paragraphs[0]
    p.text = txt
    p.font.name = "Segoe UI"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT

# Style Table Content
for r, row in enumerate(tbl_data):
    for c, val in enumerate(row):
        cell = table.cell(r + 1, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_MIDNIGHT
        p = cell.text_frame.paragraphs[0]
        p.text = val
        p.font.name = "Segoe UI"
        p.font.size = Pt(10)
        p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
        
        if c == 2:
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x22, 0xC5, 0x5E)  # Green for PASS
        else:
            p.font.color.rgb = COLOR_WHITE

# Save presentation
output_pptx = "evidence/compliance_test_report.pptx"
prs.save(output_pptx)
print(f"Successfully generated compliance presentation slide deck: {output_pptx}")
